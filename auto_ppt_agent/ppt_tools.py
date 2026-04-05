from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from auto_ppt_agent.image_service import download_image
import os


OUTPUT_FILE = "outputs/generated.pptx"


# ------------------ EXPAND CONTENT ------------------
def format_bullets(points):

    expanded = []

    for p in points:
        p = str(p).strip()

        if len(p) > 5:
            expanded.append(
                p + " — key concept explained clearly"
            )

    return expanded[:5]


# ------------------ PPT CREATION ------------------
def create_ppt(slides):

    prs = Presentation()

    os.makedirs("outputs", exist_ok=True)
    os.makedirs("temp_images", exist_ok=True)

    for i, s in enumerate(slides):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # ================= BACKGROUND =================
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(245, 248, 255)
        bg.line.fill.background()

        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # ================= HEADER =================
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(45, 120, 255)
        header.line.fill.background()

        # TITLE
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(9), Inches(0.6)
        )

        tf = title_box.text_frame
        tf.text = s["title"]

        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

        # ================= CONTENT CARD =================
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.3),
            Inches(5.8), Inches(5.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.fill.background()

        # ================= BULLETS (ONLY MAIN TEXT) =================
        bullets = format_bullets(s["points"])

        y = 1.6

        for b in bullets:

            # ICON
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.8), Inches(y),
                Inches(0.12), Inches(0.12)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = RGBColor(45, 120, 255)
            icon.line.fill.background()

            # MAIN TEXT ONLY (NO LIGHT TEXT)
            txt = slide.shapes.add_textbox(
                Inches(1.0), Inches(y - 0.05),
                Inches(5.3), Inches(0.6)
            )

            tf = txt.text_frame
            tf.word_wrap = True
            tf.text = b

            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(40, 40, 40)

            y += 1.0

        # ================= IMAGE =================
        img_path = f"temp_images/image_{i}.png"

        result = download_image(s["image"], img_path)

        if result and os.path.exists(img_path):
            slide.shapes.add_picture(
                img_path,
                Inches(6.5),
                Inches(1.5),
                Inches(3.2),
                Inches(4)
            )
        else:
            print(f"[SKIP IMAGE] Slide {i}")

    prs.save(OUTPUT_FILE)

    return OUTPUT_FILE